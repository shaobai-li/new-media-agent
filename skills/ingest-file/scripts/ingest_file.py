#!/usr/bin/env python3
"""文件录入知识库：识别并校验文件，解析为 Markdown 并保存到 cache/ 目录."""
import argparse
import json
import os
import re
import sys
import zipfile
from pathlib import Path
from typing import List

import fitz
import mammoth
from markdownify import markdownify
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

SUPPORTED_TYPES = {"pdf", "pptx", "docx", "md"}


def validate_file(file_path: str) -> dict:
    """Step 1: 文件识别与校验."""
    if not os.path.exists(file_path):
        return {
            "valid": False, "file_type": None, "size_bytes": None,
            "exists": False, "error": {"code": "FILE_NOT_FOUND", "message": "File does not exist"}
        }

    _, ext = os.path.splitext(file_path)
    file_type = ext.lower().replace(".", "")
    size_bytes = os.path.getsize(file_path)

    if file_type not in SUPPORTED_TYPES:
        return {
            "valid": False, "file_type": file_type, "size_bytes": size_bytes,
            "exists": True, "error": {"code": "UNSUPPORTED_FILE_TYPE", "message": f"Unsupported file type: {file_type}"}
        }

    if file_type == "pdf":
        try:
            with open(file_path, "rb") as f:
                if f.read(4) != b"%PDF":
                    return {
                        "valid": False, "file_type": file_type, "size_bytes": size_bytes,
                        "exists": True, "error": {"code": "FILE_CORRUPTED", "message": "Invalid PDF file"}
                    }
        except:
            return {
                "valid": False, "file_type": file_type, "size_bytes": size_bytes,
                "exists": True, "error": {"code": "FILE_CORRUPTED", "message": "Invalid PDF file"}
            }

    if file_type in {"docx", "pptx"} and not zipfile.is_zipfile(file_path):
        return {
            "valid": False, "file_type": file_type, "size_bytes": size_bytes,
            "exists": True, "error": {"code": "FILE_CORRUPTED", "message": "Invalid Office file"}
        }

    return {"valid": True, "file_type": file_type, "size_bytes": size_bytes, "exists": True, "error": None}


class DocxToMarkdown:
    """DOCX -> Markdown via mammoth -> HTML -> markdownify."""

    def __init__(self):
        self.style_map = """
            p[style-name='Heading 1'] => h1:fresh
            p[style-name='Heading 2'] => h2:fresh
            p[style-name='Heading 3'] => h3:fresh
            p[style-name='Quote'] => blockquote > p
            r[style-name='Strong'] => strong
        """
        self._image_counter = 0
        self._media_dir: Path | None = None

    def _image_handler(self, image):
        with image.open() as image_bytes:
            content = image_bytes.read()
        extension = image.content_type.split("/")[-1]
        image_filename = f"img_{self._image_counter:03d}.{extension}"
        self._image_counter += 1
        image_path = self._media_dir / image_filename
        image_path.write_bytes(content)
        return {"src": f"media/{image_filename}"}

    def convert(self, file_path: Path, output_dir: Path) -> Path:
        output_dir.mkdir(parents=True, exist_ok=True)
        self._media_dir = output_dir / "media"
        self._media_dir.mkdir(parents=True, exist_ok=True)
        self._image_counter = 0
        with open(file_path, "rb") as f:
            result = mammoth.convert_to_html(
                f,
                style_map=self.style_map,
                convert_image=mammoth.images.img_element(self._image_handler),
            )
        html_content = result.value
        markdown_content = markdownify(
            html_content,
            heading_style="ATX",
            bullets="-",
            convert=["p", "h1", "h2", "h3", "h4", "h5", "h6", "ul", "ol", "li", "b", "i", "strong", "em", "a", "img", "table", "tr", "th", "td"],
        )
        md_filename = file_path.stem + ".md"
        md_path = output_dir / md_filename
        md_path.write_text(markdown_content.strip(), encoding="utf-8")
        return md_path


class PdfToMarkdown:
    """PDF -> Markdown with PyMuPDF (text, images, tables)."""

    def __init__(self):
        self.min_image_size = 40
        self.min_image_bpc = 8
        self.text_merge_threshold = 15.0
        self._saved_images: set[int] = set()

    def convert(self, file_path: Path, output_dir: Path) -> Path:
        output_dir.mkdir(parents=True, exist_ok=True)
        media_dir = output_dir / "media"
        media_dir.mkdir(parents=True, exist_ok=True)
        self._saved_images.clear()
        all_content: list[str] = []
        doc = fitz.open(str(file_path))
        try:
            for page_num, page in enumerate(doc, 1):
                all_content.append(f"\n\n## Page {page_num}\n\n")
                elements = self._extract_page_elements(page, media_dir)
                elements.sort(key=lambda e: (e["y0"], e["x0"]))
                merged_elements = self._merge_text_blocks(elements)
                for element in merged_elements:
                    all_content.append(f"\n{element['content']}\n")
        finally:
            doc.close()
        markdown_content = "".join(all_content)
        md_filename = file_path.stem + ".md"
        md_path = output_dir / md_filename
        md_path.write_text(markdown_content.strip(), encoding="utf-8")
        return md_path

    def _extract_page_elements(self, page, media_dir: Path) -> list[dict]:
        elements: list[dict] = []
        tabs = page.find_tables()
        tab_rects = [fitz.Rect(t.bbox) for t in tabs.tables]
        tp = page.get_textpage()
        dict_content = tp.extractDICT()
        for block in dict_content["blocks"]:
            if block["type"] == 0:
                bbox = fitz.Rect(block["bbox"])
                if any(bbox.intersects(r) for r in tab_rects):
                    continue
                lines_content = []
                for line in block["lines"]:
                    span_text = "".join([s["text"] for s in line["spans"]])
                    lines_content.append(span_text.strip())
                text = "".join(lines_content)
                if text:
                    elements.append({"type": "text", "y0": bbox.y0, "x0": bbox.x0, "content": text})
        for img in page.get_images(full=True):
            xref, bpc = img[0], img[4]
            if bpc < self.min_image_bpc:
                continue
            bbox_list = page.get_image_bbox(img)
            if not bbox_list:
                continue
            bbox = bbox_list[0] if isinstance(bbox_list, list) else bbox_list
            if bbox.width < self.min_image_size or bbox.height < self.min_image_size:
                continue
            if xref not in self._saved_images:
                img_data = page.parent.extract_image(xref)
                filename = f"img_{xref}.{img_data['ext']}"
                (media_dir / filename).write_bytes(img_data["image"])
                self._saved_images.add(xref)
            else:
                ext = page.parent.extract_image(xref)["ext"]
                filename = f"img_{xref}.{ext}"
            elements.append({"type": "image", "y0": bbox.y0, "x0": bbox.x0, "content": f"\n![image](media/{filename})\n"})
        for tab in tabs.tables:
            df = tab.extract()
            if not df:
                continue
            header = "| " + " | ".join(str(x or "").replace("\n", " ") for x in df[0]) + " |"
            separator = "| " + " | ".join(["---"] * len(df[0])) + " |"
            body = ["| " + " | ".join(str(x or "").replace("\n", " ") for x in row) + " |" for row in df[1:]]
            t_bbox = fitz.Rect(tab.bbox)
            elements.append({"type": "table", "y0": t_bbox.y0, "x0": t_bbox.x0, "content": "\n" + "\n".join([header, separator] + body) + "\n"})
        return elements

    def _merge_text_blocks(self, elements: list[dict]) -> list[dict]:
        if not elements:
            return []
        merged = []
        current = elements[0]
        for next_elem in elements[1:]:
            if current["type"] == "text" and next_elem["type"] == "text" and abs(next_elem["y0"] - (current["y0"] + 10)) < self.text_merge_threshold:
                current["content"] += next_elem["content"]
            else:
                merged.append(current)
                current = next_elem
        merged.append(current)
        return merged


class PptxToMarkdown:
    """PPTX -> Markdown (slides, text, pictures, tables)."""

    def __init__(self):
        self.position_tolerance = 10.0
        self._image_counter = 0
        self._media_dir: Path | None = None

    def convert(self, file_path: Path, output_dir: Path) -> Path:
        output_dir.mkdir(parents=True, exist_ok=True)
        self._media_dir = output_dir / "media"
        self._media_dir.mkdir(parents=True, exist_ok=True)
        self._image_counter = 0
        presentation = Presentation(str(file_path))
        md_lines = [f"# {file_path.stem}\n\n", f"Total slides: {len(presentation.slides)}\n\n"]
        for slide_num, slide in enumerate(presentation.slides, 1):
            md_lines.append(f"\n---\n\n## Slide {slide_num}\n\n")
            slide_content = self._extract_slide_content(slide, slide_num)
            md_lines.extend(slide_content)
        markdown_content = "".join(md_lines)
        md_filename = file_path.stem + ".md"
        md_path = output_dir / md_filename
        md_path.write_text(markdown_content.strip(), encoding="utf-8")
        return md_path

    def _extract_slide_content(self, slide, slide_num: int) -> List[str]:
        lines: list[str] = []
        if slide.shapes.title and slide.shapes.title.text.strip():
            lines.append(f"### {slide.shapes.title.text.strip()}\n\n")
        shapes = self._collect_and_sort_shapes(slide)
        for shape in shapes:
            if shape.has_text_frame:
                lines.extend(self._extract_text_content(shape))
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                lines.append(self._extract_image(shape, slide_num))
            elif shape.has_table:
                lines.extend(self._extract_table(shape))
        return lines

    def _collect_and_sort_shapes(self, slide) -> list:
        meaningful_shapes: list[tuple[float, float, object]] = []
        for shape in slide.shapes:
            if shape.has_text_frame or shape.shape_type == MSO_SHAPE_TYPE.PICTURE or shape.has_table:
                top = shape.top.pt if hasattr(shape, "top") else 0
                left = shape.left.pt if hasattr(shape, "left") else 0
                meaningful_shapes.append((top, left, shape))
        meaningful_shapes.sort(key=lambda x: (x[0] // self.position_tolerance, x[1]))
        return [shape for _, _, shape in meaningful_shapes]

    def _extract_text_content(self, shape) -> list[str]:
        lines: list[str] = []
        for paragraph in shape.text_frame.paragraphs:
            text = "".join(run.text for run in paragraph.runs).strip()
            if not text:
                continue
            level = paragraph.level
            if level == 0:
                if re.match(r"^\d+\.\s", text) or re.match(r"^[•\-\*]\s", text):
                    lines.append(f"- {text}\n")
                else:
                    lines.append(f"{text}\n\n")
            else:
                indent = "  " * level
                lines.append(f"{indent}- {text}\n")
        return lines

    def _extract_image(self, shape, slide_num: int) -> str:
        image = shape.image
        ext = "jpg" if image.ext in ("jpeg", "jpg") else image.ext.lower()
        filename = f"slide_{slide_num:02d}_img_{self._image_counter:03d}.{ext}"
        image_path = self._media_dir / filename
        image_path.write_bytes(image.blob)
        self._image_counter += 1
        return f"![](media/{filename})\n\n"

    def _extract_table(self, shape) -> list[str]:
        lines = ["\n"]
        table = shape.table
        rows_data: list[list[str]] = []
        for row in table.rows:
            cells = []
            for cell in row.cells:
                text = cell.text_frame.text.replace("\n", " ").replace("\r", " ").strip()
                text = text.replace("|", "\\|")
                cells.append(text)
            rows_data.append(cells)
        if rows_data:
            lines.append("| " + " | ".join(rows_data[0]) + " |\n")
            lines.append("| " + " | ".join(["---"] * len(rows_data[0])) + " |\n")
            for row_data in rows_data[1:]:
                lines.append("| " + " | ".join(row_data) + " |\n")
            lines.append("\n")
        return lines


def convert_to_markdown(file_path: str, file_type: str, cache_dir: Path) -> Path:
    """Step 2: 解析为 Markdown，保存在 cache/ 目录里."""
    input_path = Path(file_path).resolve()
    output_dir = cache_dir / input_path.stem

    if file_type == "docx":
        return DocxToMarkdown().convert(input_path, output_dir)
    elif file_type == "pdf":
        return PdfToMarkdown().convert(input_path, output_dir)
    elif file_type == "pptx":
        return PptxToMarkdown().convert(input_path, output_dir)
    elif file_type == "md":
        output_dir.mkdir(parents=True, exist_ok=True)
        md_path = output_dir / (input_path.stem + ".md")
        md_path.write_text(input_path.read_text(encoding="utf-8"), encoding="utf-8")
        return md_path
    else:
        raise ValueError(f"Unsupported file type for conversion: {file_type}")


def main() -> int:
    parser = argparse.ArgumentParser(description="文件录入知识库")
    parser.add_argument("--input", required=True, help="输入文件路径")
    parser.add_argument("--cache-dir", type=Path, default=Path("cache"), help="缓存目录 (默认: cache/)")
    args = parser.parse_args()

    try:
        result = validate_file(args.input)
        if not result["valid"]:
            print(json.dumps(result), file=sys.stderr)
            return 1

        file_type = result["file_type"]
        md_path = convert_to_markdown(args.input, file_type, args.cache_dir)

        output = {
            "valid": True,
            "file_type": file_type,
            "size_bytes": result["size_bytes"],
            "markdown_path": str(md_path),
            "media_dir": str(md_path.parent / "media") if (md_path.parent / "media").exists() else None
        }
        print(json.dumps(output))
        return 0

    except Exception as e:
        print(f"Error: {e}", file=sys.stderr)
        return 1


if __name__ == "__main__":
    raise SystemExit(main())
